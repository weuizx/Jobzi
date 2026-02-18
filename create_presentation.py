#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Слайд 1: Титульный
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Разработка системы автоматизации подбора временного персонала на базе Telegram Bot API"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_text = """Выполнил: Куртяков А.
Научно-исследовательская работа
9 семестр, 2025-2026 уч. год"""
    subtitle_frame.text = subtitle_text
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for paragraph in subtitle_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # Слайд 2: Актуальность
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Актуальность исследования"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Неформальный рынок труда в России"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    points = [
        "20-25% от общей занятости (14-17 млн человек)",
        "Высокая скорость найма (1-3 дня от вакансии до выхода)",
        "Массовость (3-10 человек одновременно)",
        "Отсутствие формальных требований"
    ]

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Проблемы работодателей"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    problems = [
        "Хаос в обработке откликов (разные форматы)",
        "80% времени уходит на сбор и сортировку данных",
        "Отсутствие автоматизации",
        "Потеря истории взаимодействия с кандидатами"
    ]

    for problem in problems:
        p = tf.add_paragraph()
        p.text = problem
        p.level = 1
        p.font.size = Pt(18)

    # Слайд 3: Цели и задачи
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Цели и задачи исследования"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Цель"
    p.level = 0
    p.font.size = Pt(22)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Сократить время работодателя на подбор персонала с 3-5 часов в неделю до 10-15 минут за счет автоматизации"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Основные задачи"
    p.level = 0
    p.font.size = Pt(22)
    p.font.bold = True

    tasks = [
        "Проектирование архитектуры с мультитенантностью",
        "Реализация функционала для работодателей (вакансии, анкеты, отклики)",
        "Реализация функционала для соискателей (отклик по коду)",
        "Система рассылок в Telegram-каналы",
        "Обеспечение целостности данных (snapshot контекста)"
    ]

    for task in tasks:
        p = tf.add_paragraph()
        p.text = task
        p.level = 1
        p.font.size = Pt(16)

    # Слайд 4: Анализ существующих решений
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Анализ существующих решений"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    solutions = [
        ("HeadHunter, SuperJob", "Ориентация на долгосрочную занятость, требуют резюме, долгий цикл найма"),
        ("Avito Работа", "Отклики в разном формате, нет инструментов для анкет, ручная обработка"),
        ("YouDo, Profi.ru", "Ориентация на специалистов, а не массовый найм, высокая комиссия (15-20%)"),
        ("Telegram-чаты", "Ручное управление, отклики теряются в потоке, нет структурирования")
    ]

    for solution, problem in solutions:
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = solution
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = problem
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(128, 128, 128)

    p = tf.add_paragraph()
    p.text = "Вывод: существующие решения не закрывают потребности неформального найма"
    p.level = 0
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)

    # Слайд 5: Решение - Jobzi
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Разрабатываемое решение - Jobzi"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Telegram-бот для автоматизации подбора временного персонала"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    features = [
        "Создание вакансий за 2-3 минуты",
        "Настраиваемые анкеты (6 типов вопросов)",
        "Уникальный код вакансии формата ABC123",
        "Автоматический сбор откликов",
        "Экспорт данных в Excel",
        "Рассылка в Telegram-группы",
        "История взаимодействия с кандидатами"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1
        p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Платформа: Telegram (65+ млн пользователей в России)"
    p.level = 0
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 128, 0)

    # Слайд 6: Архитектура системы
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Архитектура системы"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    layers = [
        "Telegram API Layer - взаимодействие с Telegram Bot API",
        "Handler Layer - маршрутизация по типу пользователя",
        "Telegram Service Layer - бизнес-логика диалогов",
        "DB Service Layer - CRUD операции и бизнес-логика",
        "Repository Layer (JPA) - доступ к БД",
        "PostgreSQL Database - хранение данных (14 таблиц)"
    ]

    for layer in layers:
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = layer
        p.level = 0
        p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Многоуровневая архитектура обеспечивает слабую связанность, упрощает тестирование и масштабирование"
    p.level = 0
    p.font.size = Pt(14)
    p.font.italic = True

    # Слайд 7: Технологический стек
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Технологический стек"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    stack = [
        ("Язык программирования", "Kotlin - null safety, coroutines, совместимость с Java"),
        ("Framework", "Spring Boot - DI, JPA, транзакции, scheduling"),
        ("База данных", "PostgreSQL 16 - ACID, JSON support, constraints"),
        ("Миграции", "Liquibase - версионирование схемы БД"),
        ("Telegram Integration", "TelegramBots Spring Boot Starter"),
        ("Excel генерация", "Apache POI"),
        ("Инфраструктура", "Docker Compose, Gradle")
    ]

    for tech, description in stack:
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = tech
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = description
        p.level = 1
        p.font.size = Pt(14)

    # Слайд 8: Функционал для работодателей
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Функционал для работодателей"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Управление вакансиями"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    vacancy_features = [
        "Создание через пошаговый диалог",
        "Генерация уникального кода (ABC123)",
        "Статусы: DRAFT, ACTIVE, PAUSED, CLOSED",
        "Редактирование любого поля"
    ]

    for feature in vacancy_features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "Управление откликами"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    application_features = [
        "Просмотр всех откликов с фильтрацией",
        "Изменение статусов (NEW → VIEWED → ACCEPTED/REJECTED)",
        "Добавление заметок о кандидатах",
        "Экспорт в Excel с форматированием"
    ]

    for feature in application_features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1
        p.font.size = Pt(16)

    # Слайд 9: Функционал для соискателей
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Функционал для соискателей"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Процесс отклика"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    steps = [
        "Ввод кода вакансии (ABC123)",
        "Просмотр информации о вакансии",
        "Пошаговое заполнение анкеты",
        "Валидация ответов в реальном времени",
        "Получение подтверждения"
    ]

    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.level = 1
        p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Типы вопросов анкеты"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    question_types = [
        "TEXT - текстовый ответ",
        "PHONE - номер телефона с валидацией",
        "NUMBER - числовой ответ",
        "DATE - дата",
        "YES_NO - да/нет",
        "CHOICE - выбор из вариантов"
    ]

    for qtype in question_types:
        p = tf.add_paragraph()
        p.text = qtype
        p.level = 1
        p.font.size = Pt(16)

    # Слайд 10: Модель данных
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Модель данных"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "14 таблиц PostgreSQL"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    tables = [
        "businesses, users, business_users - мультитенантность",
        "vacancies - вакансии с уникальными кодами",
        "questions - анкеты (6 типов вопросов)",
        "applications - отклики кандидатов",
        "answers - ответы со snapshot контекста",
        "broadcast_channels, broadcast_campaigns - рассылки"
    ]

    for table in tables:
        p = tf.add_paragraph()
        p.text = table
        p.level = 1
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "Ключевые механизмы"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    mechanisms = [
        "Snapshot контекста - защита от потери данных при изменении вопросов",
        "UNIQUE constraints - защита от дублей откликов",
        "CASCADE DELETE - автоматическая очистка связанных данных",
        "Индексы - оптимизация поиска по кодам и фильтрации"
    ]

    for mechanism in mechanisms:
        p = tf.add_paragraph()
        p.text = mechanism
        p.level = 1
        p.font.size = Pt(14)

    # Слайд 11: Тестирование
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Тестирование системы"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Функциональное тестирование"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    functional_tests = [
        "Создание вакансии со стандартными вопросами",
        "Отклик с заполнением анкеты",
        "Защита от дублей откликов",
        "Валидация кодов вакансий и номеров телефонов"
    ]

    for test in functional_tests:
        p = tf.add_paragraph()
        p.text = test
        p.level = 1
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "Интеграционное тестирование"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    integration_tests = [
        "Полный цикл создания вакансии (50 итераций)",
        "Создание отклика с snapshot контекста (10 откликов)",
        "Проверка сохранения snapshot при изменении вопроса"
    ]

    for test in integration_tests:
        p = tf.add_paragraph()
        p.text = test
        p.level = 1
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "Все тесты пройдены успешно"
    p.level = 0
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    # Слайд 12: Результаты и эффективность
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Достигнутые результаты"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Сокращение временных затрат"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "До внедрения: 4-6 часов/неделю"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "После внедрения: 10-15 минут/неделю"
    p.level = 1
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Экономия: 95-97% (в 20-30 раз)"
    p.level = 1
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "Реализованный функционал"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    features = [
        "Создание вакансий за 5 шагов (~2-3 минуты)",
        "6 типов вопросов в анкетах",
        "5 статусов откликов",
        "Экспорт в Excel",
        "Система рассылок с защитой от rate limit",
        "Snapshot контекста для защиты данных"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1
        p.font.size = Pt(14)

    # Слайд 13: Дальнейшее развитие
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Дальнейшее развитие"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Подготовка к внедрению"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    preparation = [
        "Развёртывание на production (облако)",
        "Настройка резервного копирования",
        "Миграция состояний на Redis",
        "Подготовка документации"
    ]

    for item in preparation:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "Пилотное тестирование"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    pilot = [
        "Привлечение 3-5 пилотных клиентов",
        "Сопровождение первых пользователей",
        "Сбор обратной связи и метрик",
        "Выявление проблем UX"
    ]

    for item in pilot:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(16)

    # Слайд 14: Заключение
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Заключение"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Разработана система Jobzi - Telegram-бот для автоматизации подбора временного персонала"
    p.level = 0
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Ключевые достижения"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True

    achievements = [
        "Сокращение времени на подбор в 20-30 раз",
        "Полнофункциональная система с вакансиями, анкетами, откликами",
        "Успешное тестирование всех функций",
        "Готовность к пилотному внедрению"
    ]

    for achievement in achievements:
        p = tf.add_paragraph()
        p.text = achievement
        p.level = 1
        p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = "Система решает реальную проблему неформального рынка труда и готова к апробации в реальных условиях"
    p.level = 0
    p.font.size = Pt(16)
    p.font.italic = True

    # Слайд 15: Спасибо за внимание
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    thanks_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Спасибо за внимание!\n\nВопросы?"
    thanks_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for paragraph in thanks_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True

    # Сохраняем презентацию
    prs.save('9_сем_НИР_Jobzi/Презентация_НИР_Jobzi.pptx')
    print("Презентация успешно создана: 9_сем_НИР_Jobzi/Презентация_НИР_Jobzi.pptx")

def create_speaker_notes():
    notes = """ТЕКСТ ДЛЯ ВЫСТУПЛЕНИЯ НА 5-10 МИНУТ
===================================

СЛАЙД 1: ТИТУЛЬНЫЙ (10 сек)
Здравствуйте! Представляю вашему вниманию результаты научно-исследовательской работы по разработке системы автоматизации подбора временного персонала на базе Telegram Bot API.

СЛАЙД 2: АКТУАЛЬНОСТЬ (50 сек)
Начну с актуальности темы. В России существует обширный неформальный рынок труда, который составляет 20-25% от общей занятости - это около 14-17 миллионов человек. Этот сегмент характеризуется высокой скоростью найма - от публикации вакансии до выхода на работу проходит всего 1-3 дня. Также характерна массовость - нужно одновременно нанять 3-10 человек, и полное отсутствие формальных требований - не нужны резюме, дипломы, трудовые книжки.

При этом работодатели сталкиваются с серьезными проблемами. Во-первых, это хаос в обработке откликов - сообщения приходят в разных форматах и теряются в чатах. Во-вторых, до 80% времени работодателя уходит не на подбор, а на рутинную работу - сбор данных, составление списков. Также отсутствует автоматизация и теряется история взаимодействия с кандидатами.

СЛАЙД 3: ЦЕЛИ И ЗАДАЧИ (40 сек)
Цель исследования - сократить время работодателя на подбор персонала с 3-5 часов в неделю до 10-15 минут за счет автоматизации рутинных операций.

Для достижения этой цели были поставлены следующие задачи: спроектировать архитектуру системы с поддержкой мультитенантности, реализовать функционал для работодателей - создание вакансий, настройка анкет, обработка откликов, экспорт в Excel. Также реализовать функционал для соискателей - отклик по короткому коду и заполнение анкеты. Дополнительно создать систему рассылок в Telegram-каналы и обеспечить целостность данных через механизм snapshot контекста.

СЛАЙД 4: АНАЛИЗ СУЩЕСТВУЮЩИХ РЕШЕНИЙ (40 сек)
Был проведен анализ существующих решений на рынке. HeadHunter и SuperJob ориентированы на формальную занятость с долгим циклом найма. Avito Работа не предоставляет инструментов для структурированного сбора откликов. YouDo и Profi.ru фокусируются на услугах специалистов, а не на массовом наборе рабочих. Telegram-чаты, хотя и популярны, требуют полностью ручного управления, и в них отклики теряются в общем потоке сообщений.

Вывод: существующие решения не закрывают потребности сегмента временного неформального найма.

СЛАЙД 5: РЕШЕНИЕ - JOBZI (30 сек)
В рамках работы разработана система Jobzi - Telegram-бот для автоматизации подбора временного персонала. Система позволяет создавать вакансии за 2-3 минуты, использовать настраиваемые анкеты с 6 типами вопросов, генерировать уникальные короткие коды вакансий формата ABC123, автоматически собирать отклики, экспортировать данные в Excel и делать рассылки в Telegram-группы.

Выбор Telegram как платформы обусловлен тем, что им пользуется более 65 миллионов человек в России, и он уже активно используется малым бизнесом.

СЛАЙД 6: АРХИТЕКТУРА (40 сек)
Система построена по многоуровневой архитектуре. На верхнем уровне располагается Telegram API Layer для взаимодействия с ботом. Далее идет Handler Layer для маршрутизации запросов по типу пользователя - суперадмин, работодатель или соискатель. Telegram Service Layer отвечает за бизнес-логику диалогов и управление состояниями. DB Service Layer реализует CRUD операции. Repository Layer обеспечивает доступ к данным через JPA. И на нижнем уровне находится PostgreSQL с 14 таблицами.

Такая архитектура обеспечивает слабую связанность компонентов, упрощает тестирование и позволяет масштабировать систему.

СЛАЙД 7: ТЕХНОЛОГИЧЕСКИЙ СТЕК (30 сек)
В качестве языка программирования выбран Kotlin благодаря null safety, корутинам и совместимости с Java-экосистемой. Framework - Spring Boot для DI, JPA и управления транзакциями. База данных - PostgreSQL 16 с поддержкой ACID-транзакций и JSON. Для миграций используется Liquibase. Интеграция с Telegram реализована через библиотеку TelegramBots Spring Boot Starter. Для генерации Excel используется Apache POI.

СЛАЙД 8: ФУНКЦИОНАЛ ДЛЯ РАБОТОДАТЕЛЕЙ (40 сек)
Функционал для работодателей включает управление вакансиями. Создание происходит через пошаговый диалог с автоматической генерацией уникального кода. Вакансии могут находиться в статусах DRAFT, ACTIVE, PAUSED или CLOSED. Любое поле можно отредактировать.

Управление откликами позволяет просматривать все отклики с фильтрацией, менять статусы от NEW через VIEWED до ACCEPTED или REJECTED, добавлять заметки о кандидатах и экспортировать данные в Excel с форматированием.

СЛАЙД 9: ФУНКЦИОНАЛ ДЛЯ СОИСКАТЕЛЕЙ (30 сек)
Процесс отклика для соискателя максимально упрощен. Достаточно ввести код вакансии из объявления, просмотреть информацию о вакансии и пошагово заполнить анкету. Система поддерживает 6 типов вопросов: текстовый ответ, номер телефона с валидацией, число, дата, да/нет и выбор из вариантов. Все ответы валидируются в реальном времени, и соискатель получает подтверждение об отправке отклика.

СЛАЙД 10: МОДЕЛЬ ДАННЫХ (40 сек)
Модель данных включает 14 таблиц PostgreSQL. Таблицы businesses, users и business_users обеспечивают мультитенантность - изоляцию данных разных работодателей. Таблица vacancies хранит вакансии с уникальными кодами. Questions содержит анкеты с 6 типами вопросов. Applications хранит отклики, а Answers - ответы со snapshot контекста.

Ключевой механизм - snapshot контекста. При сохранении ответа в таблице answers создается копия текста вопроса, его типа и порядка. Это защищает от потери данных, если работодатель изменит или удалит вопрос после того, как соискатель уже ответил. Также используются UNIQUE constraints для защиты от дублей откликов и индексы для оптимизации поиска.

СЛАЙД 11: ТЕСТИРОВАНИЕ (30 сек)
Проведено комплексное тестирование системы. Функциональное тестирование покрыло создание вакансий, отклики с заполнением анкет, защиту от дублей и валидацию входных данных.

Интеграционное тестирование включало полный цикл создания вакансии - протестировано 50 итераций, и создание откликов со snapshot контекста - 10 откликов. Дополнительно проверена корректность работы snapshot при изменении вопросов. Все тесты пройдены успешно.

СЛАЙД 12: РЕЗУЛЬТАТЫ И ЭФФЕКТИВНОСТЬ (40 сек)
Система обеспечивает значительное сокращение временных затрат. До внедрения работодатель тратил 4-6 часов в неделю на подбор персонала. После внедрения это время сократилось до 10-15 минут в неделю. Это означает экономию в 95-97%, то есть сокращение времени в 20-30 раз.

Реализован полный функционал: создание вакансий за 5 шагов, 6 типов вопросов в анкетах, 5 статусов откликов, экспорт в Excel, система рассылок с защитой от rate limit Telegram API и механизм snapshot контекста для защиты данных.

СЛАЙД 13: ДАЛЬНЕЙШЕЕ РАЗВИТИЕ (30 сек)
Следующий этап - подготовка к внедрению. Планируется развернуть систему на production-окружении в облаке, настроить резервное копирование, мигрировать состояния диалогов на Redis для устойчивости к перезапускам и подготовить документацию.

Затем будет проведено пилотное тестирование с привлечением 3-5 клиентов из целевого сегмента. Это позволит собрать обратную связь, получить метрики использования и выявить проблемы UX в реальных условиях.

СЛАЙД 14: ЗАКЛЮЧЕНИЕ (20 сек)
В заключение: разработана система Jobzi - Telegram-бот для автоматизации подбора временного персонала. Ключевые достижения - это сокращение времени на подбор в 20-30 раз, реализация полнофункциональной системы, успешное тестирование всех функций и готовность к пилотному внедрению. Система решает реальную проблему неформального рынка труда и готова к апробации.

СЛАЙД 15: СПАСИБО ЗА ВНИМАНИЕ (5 сек)
Спасибо за внимание! Готов ответить на ваши вопросы.

===================================
ОБЩЕЕ ВРЕМЯ: ~7-8 минут (в комфортном темпе)

СОВЕТЫ ПО ВЫСТУПЛЕНИЮ:
- Поддерживайте зрительный контакт с аудиторией
- Не читайте текст слайдов дословно - они являются визуальной поддержкой
- Используйте указку или курсор для указания на важные элементы
- Делайте паузы между слайдами для вопросов
- Будьте готовы углубиться в детали, если последуют вопросы
- Держите темп - не спешите, но и не затягивайте

ВОЗМОЖНЫЕ ВОПРОСЫ И ОТВЕТЫ:

1. Почему выбрали Telegram, а не веб-приложение?
Ответ: Целевая аудитория - работодатели в неформальном секторе уже используют Telegram для коммуникации. Нулевой барьер входа, мгновенные уведомления, привычный интерфейс.

2. Как обеспечивается безопасность данных?
Ответ: Мультитенантность через business_id - каждый бизнес видит только свои данные. Валидация ввода, защита от SQL-инъекций через JPA, UNIQUE constraints в БД.

3. Как система масштабируется?
Ответ: Stateless архитектура сервисов, возможность вынесения состояний в Redis, разделение БД по бизнесам, кэширование частых запросов.

4. Что такое snapshot контекста и зачем он нужен?
Ответ: При сохранении ответа создается копия вопроса. Это защищает от потери данных - даже если работодатель изменит вопрос, в откликах останется исходная версия.

5. Какова бизнес-модель?
Ответ: Планируется подписочная модель 500-1000₽/месяц за бизнес. Целевая аудитория готова платить за экономию 4-6 часов в неделю.
"""

    with open('9_сем_НИР_Jobzi/Текст_для_выступления.txt', 'w', encoding='utf-8') as f:
        f.write(notes)

    print("Текст для выступления создан: 9_сем_НИР_Jobzi/Текст_для_выступления.txt")

if __name__ == "__main__":
    create_presentation()
    create_speaker_notes()
